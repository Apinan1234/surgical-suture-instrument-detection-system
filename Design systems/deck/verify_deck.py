# -*- coding: utf-8 -*-
"""ตรวจไฟล์ .pptx ที่สร้างแล้ว สองอย่าง

  1. ตัวเลขทุกค่าที่อ้างในสไลด์ ต้องหาเจอในไฟล์ผลจริง
  2. ข้อความในสไลด์ต้องไม่ละเมิดกติกาความถูกต้อง ❶-❾ ใน SLIDE-BRIEF.md

    python verify_deck.py
"""
import os
import re
import sys

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "..", "assets")
DECK = os.path.join(HERE, "ระบบตรวจจับอุปกรณ์การเย็บแผล-สอบโครงงาน2.pptx")

# (ค่าที่อ้างในสไลด์, ไฟล์ต้นทางที่ต้องเจอค่านี้)
NUMBERS = [
    ("0.7160", "results-960/val/aug960_overall_metrics.md"),
    ("0.4532", "results-960/val/aug960_overall_metrics.md"),
    ("0.7335", "results-960/val/aug960_overall_metrics.md"),
    ("0.7173", "results-960/val/aug960_overall_metrics.md"),
    ("0.7015", "results-960/test/aug960_overall_metrics.md"),
    ("0.4348", "results-960/test/aug960_overall_metrics.md"),
    ("0.7331", "results-960/test/aug960_overall_metrics.md"),
    ("0.6926", "results-960/test/aug960_overall_metrics.md"),
    ("0.5746", "results-comparison/val/benchmark_comparison.md"),
    ("0.6559", "results-comparison/val/benchmark_comparison.md"),
    ("0.7163", "results-comparison/val/benchmark_comparison.md"),
    ("0.3245", "results-comparison/val/benchmark_comparison.md"),
    ("0.4007", "results-comparison/val/benchmark_comparison.md"),
    ("0.4366", "results-comparison/val/benchmark_comparison.md"),
    ("0.6985", "results-comparison/val/benchmark_comparison.md"),
    ("0.6872", "results-comparison/val/benchmark_comparison.md"),
    ("0.5417", "results-comparison/val/benchmark_comparison.md"),
    ("0.6641", "results-comparison/val/benchmark_comparison.md"),
    ("0.3603", "results-960/val/aug960_per_class_metrics.md"),
    ("0.3579", "results-960/val/aug960_per_class_metrics.md"),
    ("0.8513", "results-960/val/aug960_per_class_metrics.md"),
    ("0.8662", "results-960/val/aug960_per_class_metrics.md"),
    ("0.6583", "results-960/val/aug960_per_class_metrics.md"),
    ("139.8", "results-perf/inference_benchmark.md"),
    ("155.0", "results-perf/inference_benchmark.md"),
    ("179.0", "results-perf/inference_benchmark.md"),
    ("5.59", "results-perf/inference_benchmark.md"),
    ("44.7", "results-perf/inference_benchmark.md"),
    ("57.0", "results-perf/inference_benchmark.md"),
    ("74.6", "results-perf/inference_benchmark.md"),
    ("13.41", "results-perf/inference_benchmark.md"),
    ("0.2435", "results-comparison/val/aug150_per_class_metrics.md"),
    ("0.3047", "results-comparison/val/aug150_per_class_metrics.md"),
    ("0.3351", "results-960/test/aug960_per_class_metrics.md"),
    ("0.3556", "results-960/test/aug960_per_class_metrics.md"),
    ("0.2513", "results-comparison/test/aug150_per_class_metrics.md"),
    ("0.3163", "results-comparison/test/aug150_per_class_metrics.md"),
    ("12.4", "results-boxsize/box_size_by_class.md"),
    ("80.9", "results-boxsize/box_size_by_class.md"),
    ("111.9", "results-boxsize/box_size_by_class.md"),
    ("128.7", "results-boxsize/box_size_by_class.md"),
    ("9,884", "results-960/val/class_distribution.md"),
    ("2,729", "results-960/val/class_distribution.md"),
    ("1,944", "results-boxsize/box_size_by_class.md"),
    # รอบ re-annotated — baseline9-reannot คือโมเดลหลักที่นำเสนอตั้งแต่ 16 ส.ค. 2569 รอบที่ 4
    ("0.7225", "results-reannot/val/baseline9-reannot_overall_metrics.md"),
    ("0.4444", "results-reannot/val/baseline9-reannot_overall_metrics.md"),
    ("0.7212", "results-reannot/test/baseline9-reannot_overall_metrics.md"),
    ("0.4352", "results-reannot/test/baseline9-reannot_overall_metrics.md"),
    ("0.7359", "results-reannot/val/baseline9-reannot_overall_metrics.md"),
    ("0.7290", "results-reannot/val/baseline9-reannot_overall_metrics.md"),
    ("0.7158", "results-reannot/test/baseline9-reannot_overall_metrics.md"),
    ("0.7645", "results-reannot/test/baseline9-reannot_overall_metrics.md"),
    ("0.4229", "results-reannot/val/baseline9-reannot_per_class_metrics.md"),
    ("0.3893", "results-reannot/val/baseline9-reannot_per_class_metrics.md"),
    ("0.6388", "results-reannot/val/baseline9-reannot_per_class_metrics.md"),
    ("0.9721", "results-reannot/val/baseline9-reannot_per_class_metrics.md"),
    ("0.4000", "results-reannot/test/baseline9-reannot_per_class_metrics.md"),
    ("0.3338", "results-reannot/test/baseline9-reannot_per_class_metrics.md"),
    ("0.7169", "results-reannot/screening/benchmark_comparison.md"),
    ("0.4312", "results-reannot/screening/benchmark_comparison.md"),
    ("0.7062", "results-reannot/screening/benchmark_comparison.md"),
    ("0.4272", "results-reannot/screening/benchmark_comparison.md"),
    ("0.6807", "results-reannot/screening/benchmark_comparison.md"),
    ("0.4310", "results-reannot/screening/benchmark_comparison.md"),
    ("0.5343", "results-reannot/screening/benchmark_comparison.md"),
    ("0.3006", "results-reannot/screening/benchmark_comparison.md"),
    ("12.3", "results-reannot/boxsize/box_size_by_class.md"),
    ("81.5", "results-reannot/boxsize/box_size_by_class.md"),
]

# (ชื่อกติกา, regex ที่ห้ามเจอ, คำอธิบาย)
FORBIDDEN = [
    ("❸", r"0\.608", "ตัวเลข 0.608 ของโมเดล 5 คลาสเดิม ทำซ้ำไม่ได้แล้ว ห้ามอ้าง"),
    ("❶", r"ความละเอียด[^\n]{0,40}(คะแนนรวม|mAP50)[^\n]{0,20}(ดีขึ้น|สูงขึ้น|เพิ่มขึ้น)",
     "ห้ามเคลมว่าเพิ่มความละเอียดแล้วคะแนนรวมดีขึ้น"),
    ("⓫", r"\bv2\b[^\n]{0,60}reannot|reannot[^\n]{0,60}\bv2\b",
     "ห้ามใช้คำว่า v2 คู่กับ reannot — ใช้ reannot เท่านั้น"),
]


def slide_text(slide):
    out = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            out.append(shp.text_frame.text)
        if shp.has_table:
            for row in shp.table.rows:
                out.extend(c.text for c in row.cells)
    return "\n".join(out)


def main():
    prs = Presentation(DECK)
    texts = [slide_text(s) for s in prs.slides]
    everything = "\n".join(texts)
    fails = 0

    print("=== 1. ตัวเลขในสไลด์ เทียบกับไฟล์ต้นทาง ===")
    cache = {}
    for value, rel in NUMBERS:
        path = os.path.join(ASSETS, rel)
        if rel not in cache:
            cache[rel] = open(path, encoding="utf-8").read()
        in_deck = value in everything
        in_src = value in cache[rel] or value.replace(",", "") in cache[rel]
        if not in_deck:
            print("  ข้าม  %-9s ไม่ได้ใช้ในสไลด์" % value)
        elif not in_src:
            print("  ผิด   %-9s อ้างในสไลด์ แต่ไม่พบใน %s" % (value, rel))
            fails += 1
    print("  ตรวจ %d ค่า" % len(NUMBERS))

    print("\n=== 2. กติกาความถูกต้อง ===")
    for rule, pat, why in FORBIDDEN:
        hits = []
        for i, t in enumerate(texts, 1):
            if re.search(pat, t):
                hits.append(i)
        if hits:
            print("  ผิดกติกา %s  สไลด์ %s — %s" % (rule, hits, why))
            fails += 1
        else:
            print("  ผ่าน กติกา %s — %s" % (rule, why))

    # ❷ ทุกสไลด์ที่พูดถึง 0.7160 ต้องมีหมายเหตุ Stitch Scissors อยู่ในสไลด์เดียวกัน
    bad2 = [i for i, t in enumerate(texts, 1)
            if "0.7160" in t and "Stitch Scissors" not in t and "0.6990" not in t]
    if bad2:
        print("  ผิดกติกา ❷  สไลด์ %s พูดถึง 0.7160 โดยไม่มีหมายเหตุ Stitch Scissors" % bad2)
        fails += 1
    else:
        print("  ผ่าน กติกา ❷ — ทุกสไลด์ที่อ้าง 0.7160 มีหมายเหตุ Stitch Scissors ในแผ่นเดียวกัน")

    # ❹ ตัวเลขความเร็วพูดได้ตั้งแต่ 14 ส.ค. 2569 แต่ต้องอ้างไฟล์ผลการวัดในสไลด์เดียวกัน
    perf = [i for i, t in enumerate(texts, 1)
            if re.search(r"\bms\b|เฟรม/วินาที|เฟรม/นาที|\bFPS\b", t)]
    bad4 = [i for i in perf if "inference_benchmark.md" not in texts[i - 1]]
    if bad4:
        print("  ผิดกติกา ❹  สไลด์ %s อ้างตัวเลขความเร็วโดยไม่อ้างไฟล์ผลการวัด" % bad4)
        fails += 1
    else:
        print("  ผ่าน กติกา ❹ — ตัวเลขความเร็วอยู่ใน %d สไลด์ ทุกแผ่นอ้าง inference_benchmark.md" % len(perf))

    # ❿ ห้ามเขียนว่ารอบ reannot "ชนะ"/"ดีกว่า" aug960 — คนละ val split เทียบกันตรง ๆ ไม่ได้
    bad10 = [i for i, t in enumerate(texts, 1)
             if re.search(r"reannot|baseline9-reannot|0\.7225|0\.7212", t)
             and re.search(r"ชนะ|ดีกว่า.{0,20}aug960|aug960.{0,20}ดีกว่า", t)]
    if bad10:
        print("  ผิดกติกา ⓮/❿  สไลด์ %s เขียนว่ารอบ reannot ชนะ/ดีกว่า aug960 ทั้งที่คนละ val split" % bad10)
        fails += 1
    else:
        print("  ผ่าน กติกา ⓮/❿ — ไม่มีสไลด์ไหนเคลมว่ารอบ reannot ชนะ/ดีกว่า aug960 ทางสถิติ")

    # W2: สไลด์คั่นบทของส่วน 2/3/4/5/6 ต้องอ้างกลับไปที่ผัง W-1 (บอกว่ากำลังขยายขั้นไหน)
    _SECTION_TITLES = (r"^Model$", r"^Data Collection$", r"^Model Training$", r"^Model Result$",
                        r"^Web Application Demonstration$")
    section_dividers = {i: t for i, t in enumerate(texts, 1)
                         if any(re.search(pat, line.strip())
                                for line in t.split("\n") for pat in _SECTION_TITLES)}
    missing_w2 = [i for i, t in section_dividers.items() if "ขยายขั้นที่" not in t]
    if section_dividers and missing_w2:
        print("  ผิดกติกา W2  สไลด์คั่นบท %s ไม่มีคำว่า \"ขยายขั้นที่\" อ้างกลับไปผัง W-1" % missing_w2)
        fails += 1
    elif section_dividers:
        print("  ผ่าน กติกา W2 — สไลด์คั่นบทของส่วน 2/3/4/5/6 ทั้ง %d แผ่นอ้างกลับไปผัง W-1 ครบ"
              % len(section_dividers))
    else:
        print("  เตือน  หาสไลด์คั่นบทของส่วน 2/3/4/5/6 ไม่เจอเลย — ตรวจ regex ใน verify_deck.py")

    # ทุกสไลด์ที่มีตัวเลขทศนิยม ต้องมีบรรทัดที่มา
    no_src = [i for i, t in enumerate(texts, 1)
              if re.search(r"\d\.\d{3,4}", t) and "ที่มา:" not in t]
    if no_src:
        print("  เตือน  สไลด์ %s มีตัวเลขแต่ไม่มีบรรทัด 'ที่มา:'" % no_src)
    else:
        print("  ผ่าน  ทุกสไลด์ที่มีตัวเลขมีบรรทัดที่มากำกับ")

    print("\nสรุป: %s (%d สไลด์)" % ("ไม่พบปัญหา" if not fails else "พบปัญหา %d จุด" % fails,
                                    len(texts)))
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
