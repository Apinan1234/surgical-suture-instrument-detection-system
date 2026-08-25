# -*- coding: utf-8 -*-
"""สร้างไฟล์ .pptx จาก content.py

    python build_deck.py

สิ่งที่สคริปต์นี้ทำนอกจากวางสไลด์:
  1. ยัด design system ลง theme1.xml ของไฟล์ — สีทั้ง 7 บทบาทจะไปโผล่ในแถบเลือกสีของ
     PowerPoint และฟอนต์ Leelawadee UI จะเป็นฟอนต์ธีม (รวม <a:cs> สำหรับภาษาไทย)
     ทำให้กล่องข้อความที่ผู้ใช้เพิ่มเองทีหลังยังอยู่ในชุดเดียวกัน
  2. ตรวจว่าไม่มีวัตถุใดล้นขอบสไลด์ แล้วรายงานออกมา
"""
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import theme as T          # noqa: E402
import layouts as L        # noqa: E402
import content             # noqa: E402

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "ระบบตรวจจับอุปกรณ์การเย็บแผล-สอบโครงงาน2.pptx")

_CLR_ORDER = ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
              "accent4", "accent5", "accent6", "hlink", "folHlink"]


def inject_theme(prs):
    """เขียนสีและฟอนต์ของ design system ลง theme1.xml ให้ PowerPoint รู้จัก"""
    part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    xml = part.blob.decode("utf-8")

    clr = ['<a:clrScheme name="ระบบตรวจจับอุปกรณ์การเย็บแผล">']
    for key in _CLR_ORDER:
        clr.append('<a:%s><a:srgbClr val="%s"/></a:%s>' % (key, T.THEME_COLORS[key], key))
    clr.append("</a:clrScheme>")
    xml, n1 = re.subn(r"<a:clrScheme.*?</a:clrScheme>", "".join(clr), xml, flags=re.S)

    def font_block(tag):
        return ('<a:%s><a:latin typeface="%s"/><a:ea typeface="%s"/>'
                '<a:cs typeface="%s"/></a:%s>' % (tag, T.FONT, T.FONT, T.FONT, tag))

    fonts = ('<a:fontScheme name="%s">%s%s</a:fontScheme>'
             % (T.FONT, font_block("majorFont"), font_block("minorFont")))
    xml, n2 = re.subn(r"<a:fontScheme.*?</a:fontScheme>", fonts, xml, flags=re.S)

    if not (n1 == 1 and n2 == 1):
        raise RuntimeError("แทนที่ theme ไม่สำเร็จ (clrScheme=%d fontScheme=%d)" % (n1, n2))
    part._blob = xml.encode("utf-8")


def check_bounds(prs, tol_in=0.02):
    """หาวัตถุที่ล้นขอบสไลด์ — จับรูปที่คำนวณขนาดผิดได้ตั้งแต่ยังไม่ต้องเปิด PowerPoint"""
    tol = Inches(tol_in)
    sw, sh = prs.slide_width, prs.slide_height
    bad = []
    for i, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            if shp.left is None or shp.width is None:
                continue
            over = []
            if shp.left < -tol:
                over.append("ซ้าย %.2f\"" % (Emu(shp.left).inches))
            if shp.top < -tol:
                over.append("บน %.2f\"" % (Emu(shp.top).inches))
            if shp.left + shp.width > sw + tol:
                over.append("ขวา +%.2f\"" % (Emu(shp.left + shp.width - sw).inches))
            if shp.top + shp.height > sh + tol:
                over.append("ล่าง +%.2f\"" % (Emu(shp.top + shp.height - sh).inches))
            if over:
                bad.append((i, shp.shape_type, shp.name, ", ".join(over)))
    return bad


def main():
    prs = Presentation()
    prs.slide_width = Inches(T.SLIDE_W)
    prs.slide_height = Inches(T.SLIDE_H)
    inject_theme(prs)

    content.build(prs)

    total = 0
    for i, slide in enumerate(prs.slides, 1):
        total = i
        if i > 1:                       # ไม่ใส่เลขหน้าบนปก
            L.page_number(slide, i)

    prs.save(OUT)
    print("บันทึกแล้ว: %s" % OUT)
    print("จำนวนสไลด์: %d  ·  ขนาดไฟล์: %.1f MB" % (total, os.path.getsize(OUT) / 1e6))

    bad = check_bounds(prs)
    if bad:
        print("\n!! วัตถุล้นขอบสไลด์ %d จุด:" % len(bad))
        for s, st, name, why in bad:
            print("   สไลด์ %-3d %-22s %s" % (s, name, why))
    else:
        print("ตรวจขอบสไลด์: ไม่มีวัตถุล้นกรอบ")


if __name__ == "__main__":
    main()
