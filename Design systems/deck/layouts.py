# -*- coding: utf-8 -*-
"""Archetype ของสไลด์ตาม SLIDE-BRIEF.md ส่วนที่ 3 — หนึ่ง archetype = หนึ่งฟังก์ชัน

กับดักที่จัดการไว้ในไฟล์นี้:
  - ฟอนต์ไทย: python-pptx เซ็ตให้เฉพาะ <a:latin> ถ้าไม่เซ็ต <a:cs> ด้วย PowerPoint จะไปหยิบ
    ฟอนต์อื่นมาเรนเดอร์ข้อความไทย วรรณยุกต์ลอย ดู _set_cs()
  - รูปต้องขยาย/ย่อเท่ากันทุกแกนเสมอ (กฎในบรีฟ) ดู fit()
  - ทุกสไลด์ที่มีตัวเลขต้องมีที่มา: source= เป็นพารามิเตอร์ของทุก archetype ที่แสดงตัวเลขได้
"""
import os
import re

from PIL import Image as PILImage
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

import theme as T

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")
_A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'


# -- ฟอนต์ --------------------------------------------------------------------
def _set_cs(font, name):
    """เซ็ต <a:cs> (complex script) ให้ข้อความไทยเรนเดอร์ด้วยฟอนต์เดียวกับ latin"""
    rPr = font._element
    succ = ("a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = parse_xml("<a:ea %s/>" % _A)
        rPr.insert_element_before(ea, "a:cs", *succ)
    ea.set("typeface", name)
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = parse_xml("<a:cs %s/>" % _A)
        rPr.insert_element_before(cs, *succ)
    cs.set("typeface", name)


def style(font, size=None, bold=None, color=None, name=None, italic=None):
    name = name or T.FONT
    font.name = name
    _set_cs(font, name)
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic
    if color is not None:
        font.color.rgb = color
    return font


# -- ข้อความ ------------------------------------------------------------------
def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _runs(p, text, size, color, bold=False):
    """รองรับ **ตัวหนา** และ `โค้ด` แบบมาร์กดาวน์ย่อในบรรทัดเดียว"""
    on_dark = (color == T.WHITE)
    for part in re.split(r"(\*\*.+?\*\*|`[^`]+`)", text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            r = p.add_run()
            r.text = part[2:-2]
            style(r.font, size, True, T.WHITE if on_dark else T.INK)
        elif part.startswith("`") and part.endswith("`"):
            r = p.add_run()
            r.text = part[1:-1]
            style(r.font, size - 1, False, T.WHITE if on_dark else T.MUTED, name=T.FONT_MONO)
        else:
            r = p.add_run()
            r.text = part
            style(r.font, size, bold, color)


def para(tf, text, size=T.SZ_BODY, color=T.INK, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, first=False, line=1.25):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    if text:
        _runs(p, text, size, color, bold)
    return p


def bullet_list(tf, items, size=T.SZ_BODY, color=T.INK, gap=9):
    """str ธรรมดา = ระดับ 1 · ขึ้นต้นด้วยเว้นวรรค 2 ตัว = ระดับ 2 · '' = เว้นบรรทัด"""
    first = True
    for it in items:
        if it == "":
            para(tf, "", size=6, first=first, space_after=0)
            first = False
            continue
        lvl = 1 if it.startswith("  ") else 0
        marker = "▪  " if lvl == 0 else "–  "
        p = para(tf, marker + it.strip(),
                 size=size if lvl == 0 else size - 1,
                 color=color if lvl == 0 else T.MUTED,
                 first=first, space_after=gap)
        if lvl:
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Inches(0.32))))
            pPr.set("indent", "0")
        first = False


# -- ชิ้นส่วนร่วมของทุกสไลด์ ---------------------------------------------------
def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None, radius=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if radius:
        shp.adjustments[0] = radius
    shp.text_frame.word_wrap = True
    return shp


def title_bar(slide, title, eyebrow=None, dark=False):
    ink = T.WHITE if dark else T.INK
    if eyebrow:
        tf = textbox(slide, T.MARGIN, T.EYEBROW_Y, T.CONTENT_W, 0.26)
        para(tf, eyebrow, size=T.SZ_EYEBROW, color=T.ACCENT, bold=True, first=True, space_after=0)
    tf = textbox(slide, T.MARGIN, T.TITLE_Y, T.CONTENT_W, T.TITLE_H)
    para(tf, title, size=T.SZ_TITLE, color=ink, bold=True, first=True, space_after=0, line=1.1)
    rect(slide, T.MARGIN, T.RULE_Y, 1.45, 0.055, T.ACCENT)


def source_note(slide, text, warn=None):
    if not text and not warn:
        return
    tf = textbox(slide, T.MARGIN, T.SOURCE_Y, T.CONTENT_W - 1.4, 0.36)
    if warn:
        para(tf, warn, size=T.SZ_NOTE, color=T.ORANGE, first=True, space_after=2)
        if text:
            para(tf, "ที่มา: " + text, size=T.SZ_SOURCE, color=T.MUTED, space_after=0)
    else:
        para(tf, "ที่มา: " + text, size=T.SZ_SOURCE, color=T.MUTED, first=True, space_after=0)


def page_number(slide, label):
    tf = textbox(slide, T.SLIDE_W - T.MARGIN - 1.3, T.SOURCE_Y + 0.04, 1.3, 0.26)
    para(tf, str(label), size=T.SZ_SOURCE, color=T.MUTED, align=PP_ALIGN.RIGHT,
         first=True, space_after=0)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def new(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, T.INK if dark else T.WHITE)
    return s


# -- รูป ----------------------------------------------------------------------
def fit(slide, path, bx, by, bw, bh):
    """วางรูปกลางกรอบ ขยาย/ย่อเท่ากันทุกแกน (กฎในบรีฟ: ห้าม re-scale ไม่เท่ากัน)"""
    full = path if os.path.isabs(path) else os.path.join(ASSETS, path)
    with PILImage.open(full) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih)
    w, h = iw * scale, ih * scale
    return slide.shapes.add_picture(full, Inches(bx + (bw - w) / 2), Inches(by + (bh - h) / 2),
                                    Inches(w), Inches(h))


def _send_to_back(slide, pic):
    slide.shapes._spTree.remove(pic.element)
    slide.shapes._spTree.insert(2, pic.element)


def _translucent(slide, x, y, w, h, color, alpha_pct):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    srgb = shp.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (_A, int(alpha_pct * 1000))))
    return shp


# -- ARCHETYPE 1: ปก -----------------------------------------------------------
def cover(prs, title, subtitle, authors, meta, bg_image=None):
    s = new(prs, dark=True)
    if bg_image:
        pic = fit(s, bg_image, 3.2, 1.1, 9.9, 5.5)
        _send_to_back(s, pic)
        _translucent(s, 0, 0, T.SLIDE_W, T.SLIDE_H, T.INK, 88)
    rect(s, 0, 0, 0.28, T.SLIDE_H, T.ACCENT)
    tf = textbox(s, 1.15, 2.0, 10.7, 2.5)
    para(tf, title, size=42, color=T.WHITE, bold=True, first=True, space_after=14, line=1.15)
    para(tf, subtitle, size=17, color=T.ACCENT, space_after=0)
    tf2 = textbox(s, 1.15, 4.8, 10.7, 1.6)
    para(tf2, authors, size=15, color=T.WHITE, first=True, space_after=8)
    para(tf2, meta, size=12, color=T.MUTED, space_after=0)
    return s


# -- ARCHETYPE 2: คั่นบท -------------------------------------------------------
def section(prs, number, title, lead=None):
    s = new(prs, dark=True)
    rect(s, 0, 0, 0.28, T.SLIDE_H, T.ACCENT)
    tf = textbox(s, 1.35, 2.7, 10.5, 2.2)
    para(tf, number, size=13, color=T.ACCENT, bold=True, first=True, space_after=10)
    para(tf, title, size=36, color=T.WHITE, bold=True, space_after=12, line=1.15)
    if lead:
        para(tf, lead, size=15, color=T.MUTED, space_after=0)
    return s


# -- ARCHETYPE 3: ข้อความ / bullet ---------------------------------------------
def bullets(prs, title, items, eyebrow=None, source=None, warn=None, lead=None, size=T.SZ_BODY):
    s = new(prs)
    title_bar(s, title, eyebrow)
    y = T.BODY_Y
    if lead:
        tf = textbox(s, T.MARGIN, y, T.CONTENT_W, 0.75)
        para(tf, lead, size=T.SZ_SUBTITLE, color=T.NAVY2, bold=True, first=True,
             space_after=0, line=1.2)
        y += 0.9
    tf = textbox(s, T.MARGIN, y, T.CONTENT_W, T.BODY_Y + T.BODY_H - y)
    bullet_list(tf, items, size=size)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 4: diagram เต็มสไลด์ --------------------------------------------
def diagram(prs, title, image, caption=None, eyebrow=None, source=None, warn=None):
    s = new(prs)
    title_bar(s, title, eyebrow)
    cap_h = 0.5 if caption else 0.0
    fit(s, image, T.MARGIN, T.BODY_Y, T.CONTENT_W, T.BODY_H - cap_h - 0.08)
    if caption:
        tf = textbox(s, T.MARGIN, T.BODY_Y + T.BODY_H - cap_h, T.CONTENT_W, cap_h)
        para(tf, caption, size=T.SZ_BODY_SM, color=T.NAVY2, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 5: diagram ทรงสูง (รูปซ้าย ข้อความขวา) --------------------------
def tall_diagram(prs, title, image, points, eyebrow=None, source=None, warn=None, img_w=4.9):
    s = new(prs)
    title_bar(s, title, eyebrow)
    fit(s, image, T.MARGIN, T.BODY_Y, img_w, T.BODY_H)
    tx = T.MARGIN + img_w + 0.55
    tf = textbox(s, tx, T.BODY_Y + 0.1, T.SLIDE_W - tx - T.MARGIN, T.BODY_H)
    bullet_list(tf, points, size=T.SZ_BODY, gap=11)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 6: แถวตัวเลข ----------------------------------------------------
def metric_row(prs, title, metrics, eyebrow=None, note=None, source=None, warn=None, lead=None):
    """metrics: list ของ (ค่า, ป้าย, สี)"""
    s = new(prs)
    title_bar(s, title, eyebrow)
    y = T.BODY_Y
    if lead:
        tf = textbox(s, T.MARGIN, y, T.CONTENT_W, 0.6)
        para(tf, lead, size=T.SZ_SUBTITLE, color=T.NAVY2, bold=True, first=True,
             space_after=0, line=1.2)
        y += 0.85
    n = len(metrics)
    gap = 0.28
    cw = (T.CONTENT_W - gap * (n - 1)) / n
    ch = 1.95
    for i, (val, lab, col) in enumerate(metrics):
        x = T.MARGIN + i * (cw + gap)
        rect(s, x, y, cw, ch, T.SURFACE, radius=0.06)
        tf = textbox(s, x + 0.12, y + 0.4, cw - 0.24, 1.0, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, val, size=T.SZ_METRIC, color=col, bold=True, first=True, space_after=0,
             align=PP_ALIGN.CENTER, line=1.0)
        tf2 = textbox(s, x + 0.12, y + 1.38, cw - 0.24, 0.5)
        para(tf2, lab, size=T.SZ_NOTE, color=T.MUTED, first=True, space_after=0,
             align=PP_ALIGN.CENTER, line=1.15)
    if note:
        tf = textbox(s, T.MARGIN, y + ch + 0.35, T.CONTENT_W, T.BODY_Y + T.BODY_H - y - ch - 0.35)
        bullet_list(tf, note, size=T.SZ_BODY_SM, gap=8)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 7: เทียบซ้าย/ขวา ------------------------------------------------
def two_up(prs, title, left, right, eyebrow=None, source=None, warn=None, verdict=None):
    """left/right: dict(head=, sub=, items=[], accent=สี)"""
    s = new(prs)
    title_bar(s, title, eyebrow)
    cw = (T.CONTENT_W - 0.45) / 2
    h = T.BODY_H - (0.75 if verdict else 0.0)
    for i, side in enumerate((left, right)):
        x = T.MARGIN + i * (cw + 0.45)
        col = side.get("accent", T.NAVY2)
        rect(s, x, T.BODY_Y, cw, h, T.SURFACE, radius=0.04)
        rect(s, x, T.BODY_Y, cw, 0.055, col)
        tf = textbox(s, x + 0.32, T.BODY_Y + 0.3, cw - 0.64, 0.9)
        para(tf, side["head"], size=T.SZ_SUBTITLE, color=col, bold=True, first=True, space_after=4)
        if side.get("sub"):
            para(tf, side["sub"], size=T.SZ_NOTE, color=T.MUTED, space_after=0)
        tf2 = textbox(s, x + 0.32, T.BODY_Y + 1.3, cw - 0.64, h - 1.55)
        bullet_list(tf2, side["items"], size=T.SZ_BODY_SM, gap=9)
    if verdict:
        tf = textbox(s, T.MARGIN, T.BODY_Y + h + 0.2, T.CONTENT_W, 0.5)
        para(tf, verdict, size=T.SZ_BODY, color=T.INK, bold=True, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 8: ตัวเลขเดียวเต็มจอ --------------------------------------------
def big_number(prs, number, label, sub=None, eyebrow=None, source=None, color=None):
    s = new(prs, dark=True)
    rect(s, 0, 0, 0.28, T.SLIDE_H, T.ACCENT)
    if eyebrow:
        tf = textbox(s, 1.35, 1.7, 10.5, 0.4)
        para(tf, eyebrow, size=T.SZ_EYEBROW, color=T.ACCENT, bold=True, first=True, space_after=0)
    tf = textbox(s, 1.35, 2.2, 10.5, 1.8)
    para(tf, number, size=T.SZ_NUMBER, color=color or T.WHITE, bold=True, first=True,
         space_after=8, line=1.0)
    tf2 = textbox(s, 1.35, 4.2, 10.5, 2.0)
    para(tf2, label, size=T.SZ_SUBTITLE, color=T.WHITE, bold=True, first=True,
         space_after=10, line=1.25)
    if sub:
        para(tf2, sub, size=T.SZ_BODY_SM, color=T.MUTED, space_after=0, line=1.35)
    if source:
        tf3 = textbox(s, 1.35, T.SOURCE_Y, 10.5, 0.3)
        para(tf3, "ที่มา: " + source, size=T.SZ_SOURCE, color=T.MUTED, first=True, space_after=0)
    return s


# -- ARCHETYPE 9: ตาราง --------------------------------------------------------
def table(prs, title, headers, rows, col_w=None, eyebrow=None, source=None, warn=None,
          highlight=None, note=None, lead=None, size=T.SZ_TABLE):
    """highlight: set ของ index แถว (0-based ในส่วน body) ที่จะย้อมสีเน้น"""
    s = new(prs)
    title_bar(s, title, eyebrow)
    y = T.BODY_Y
    if lead:
        tf = textbox(s, T.MARGIN, y, T.CONTENT_W, 0.5)
        para(tf, lead, size=T.SZ_BODY, color=T.NAVY2, bold=True, first=True, space_after=0)
        y += 0.62
    note_lines = len(note) if note else 0
    note_h = 0.3 + 0.26 * note_lines if note else 0.0
    avail_h = T.BODY_Y + T.BODY_H - y - note_h
    nrows = len(rows) + 1
    row_h = min(0.42, avail_h / nrows)
    tbl_h = row_h * nrows
    gf = s.shapes.add_table(nrows, len(headers), Inches(T.MARGIN), Inches(y),
                            Inches(T.CONTENT_W), Inches(tbl_h))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_w:
        total = float(sum(col_w))
        for i, cv in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(T.CONTENT_W) * cv / total))
    for r in range(nrows):
        tbl.rows[r].height = Inches(row_h)
    highlight = highlight or set()

    def _cell(c, text, bold, color, fill, align):
        c.fill.solid()
        c.fill.fore_color.rgb = fill
        c.margin_left = c.margin_right = Inches(0.09)
        c.margin_top = c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.text_frame.word_wrap = True
        p = c.text_frame.paragraphs[0]
        p.alignment = align
        _runs(p, str(text), size, color, bold)

    for j, htxt in enumerate(headers):
        _cell(tbl.cell(0, j), htxt, True, T.WHITE, T.NAVY2,
              PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        base = T.SURFACE if i % 2 == 0 else T.WHITE
        fill = RGBColor(0xFD, 0xEA, 0xEE) if i in highlight else base
        for j, val in enumerate(row):
            _cell(tbl.cell(i + 1, j), val, (i in highlight and j == 0), T.INK, fill,
                  PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    if note:
        tf = textbox(s, T.MARGIN, y + tbl_h + 0.16, T.CONTENT_W, note_h)
        bullet_list(tf, note, size=T.SZ_BODY_SM, gap=6)
    source_note(s, source, warn)
    return s


# -- ARCHETYPE 10: กริดภาพหน้าจอ -----------------------------------------------
def screenshot_grid(prs, title, shots, eyebrow=None, source=None, cols=3, warn=None):
    """shots: list ของ (path, caption)"""
    s = new(prs)
    title_bar(s, title, eyebrow)
    rows = (len(shots) + cols - 1) // cols
    gap = 0.26
    cw = (T.CONTENT_W - gap * (cols - 1)) / cols
    ch = (T.BODY_H - gap * (rows - 1)) / rows
    for i, (path, cap) in enumerate(shots):
        r, c = divmod(i, cols)
        x = T.MARGIN + c * (cw + gap)
        y = T.BODY_Y + r * (ch + gap)
        rect(s, x, y, cw, ch, T.SURFACE, line=T.HAIRLINE)
        fit(s, path, x + 0.07, y + 0.07, cw - 0.14, ch - 0.46)
        tf = textbox(s, x + 0.08, y + ch - 0.34, cw - 0.16, 0.3)
        para(tf, cap, size=T.SZ_SOURCE, color=T.MUTED, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
    source_note(s, source, warn)
    return s
